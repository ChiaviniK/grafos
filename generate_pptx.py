from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()
    
    # helper for adding a slide with title and content
    def add_slide(title_text, subtitle_text, content_list, bg_image=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Subtitle or main text
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = subtitle_text
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            
        if bg_image and os.path.exists(bg_image):
            # We can't easily set as 'background' in simple python-pptx without more code
            # so we just add it as a picture in a corner or side
            prs.slides[len(prs.slides)-1].shapes.add_picture(bg_image, Inches(6), Inches(1), width=Inches(3.5))

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "IA e Data Centers: O Preço Energético da Inteligência"
    subtitle.text = "Onde os bits encontram os watts\nLuiz Chiavini - 2026"
    
    base_path = "public/assets/generated/"
    
    # Slide 2: A Explosão da Demanda
    add_slide("A Explosão da Demanda", 
              "O salto energético da IA Generativa", 
              ["Busca Google: 0.3 Wh", "Prompt ChatGPT: 2.9 Wh", "A IA consome 10x mais energia por interação."],
              os.path.join(base_path, "gpu_energy.png"))
    
    # Slide 3: Por que a IA 'bebe' tanta energia?
    add_slide("Por que a IA 'bebe' tanta energia?", 
              "Hardware e Ciclos de Treinamento", 
              ["NVIDIA H100: Consumo de até 700W (pico).", "Treino de um LLM: Energia de centenas de casas brasileiras/ano.", "Inferência constante por bilhões de usuários."],
              os.path.join(base_path, "gpu_energy.png"))
    
    # Slide 4: Anatomia do Gasto Energético
    add_slide("Anatomia do Gasto Energético", 
              "Processamento vs Refrigeração", 
              ["Processamento (GPUs): ~60% do gasto.", "Refrigeração (Cooling): ~40% para evitar superaquecimento.", "PUE: Power Usage Effectiveness (razão total/TI)."],
              os.path.join(base_path, "server_nature.png"))
    
    # Slide 5: Impacto Ambiental e Hídrico
    add_slide("Impacto Ambiental e Hídrico", 
              "Água e Geopolítica", 
              ["Consumo de bilhões de litros de água para resfriamento.", "Migração para regiões frias (Islândia, Nórdicos).", "Pressão sobre grades elétricas nacionais."],
              os.path.join(base_path, "world_map.png"))
    
    # Slide 6: Soluções: Caminho para Sustentabilidade
    add_slide("Soluções: Caminho para Sustentabilidade", 
              "Energia Limpa e Inovação", 
              ["Energias 24/7: SMRs (Nucleares Modulares) e Geotérmica.", "IA para IA: Algoritmos otimizando o refrigeração (DeepMind).", "Hardware Especializado (LPUs, TPUs, Edge Computing)."],
              os.path.join(base_path, "smr_reactor.png"))
    
    # Slide 7: O Dilema Ético e o Futuro
    add_slide("O Dilema Ético e o Futuro", 
              "Paradoxo de Jevons", 
              ["Eficiência gera mais demanda?", "Conflito entre progresso da IA e metas de emissão zero.", "Inovação responsável."],
              os.path.join(base_path, "brain_balance.png"))
    
    # Slide 8: Conclusão
    add_slide("Conclusão", 
              "Obrigado!", 
              ["'A inteligência do futuro depende da sustentabilidade do presente.'", "Perguntas?"],
              os.path.join(base_path, "server_nature.png"))

    prs.save(os.path.join(base_path, "Palestra_IA_Energia_V2.pptx"))
    print("PPTX generated successfully!")

if __name__ == "__main__":
    create_presentation()
